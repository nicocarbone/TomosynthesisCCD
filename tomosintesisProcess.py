#%%
# General imports
import matplotlib.pyplot as plt
import numpy as np
import gc
from scipy.optimize import curve_fit
import imageio
import progressbar
import importlib

# Import Tomosintesis functions
import tomosintesisFuncs as tomo
from FftGauss import fft_gauss 

# Import theory functions
import theoryContini as theo

#Scikit
from skimage.measure import compare_ssim

# %%
# Load input images

# Exp
# path ="../LAB15052019/"
# path_img = path + "TIFFs/"
# path_output = path
# images_mc, sourcesRaw = tomo.load_imagesTIFF_Zpath(path_img, "tomosyn_*.tif")
# d = 5 #Bulk thickness
# px_per_cm = 32 #Pixel per cm in input images
# source_shape = [52, 47] #Source array configuration
# expData=True


# MC
path ="../MC25112019/"
path_img = path
images_mc, sources = tomo.load_imagesMC_wsp(path_img, "*_Trans.dat", 2, posx = 4, posy = 5)
d = 5 #Bulk thickness
px_per_cm = 32 #Pixel per cm in input images
source_shape = [53, 46]
expData = False

im_shape = np.shape(images_mc[12])


# In experimental data, calculate step size from data. Else, use given 
# (as it should be precisely known) 
if expData:
    stepSizeCalc = abs((np.asarray(tomo.max_pos_fit(images_mc[0]))-
                        np.asarray(tomo.max_pos_fit(images_mc[39])))/(40*px_per_cm))[0]
    stepSize=[stepSizeCalc,stepSizeCalc]
    initPos=abs([0,im_shape[1]]-np.asarray(tomo.max_pos_fit(images_mc[0])))/px_per_cm
else:
    stepSize = [0.2, 0.2]
    initPos=[2.75, 1.75]

# Calculete max and min of images
maximg = np.amax(images_mc)
minimg = np.amin(images_mc)

# If not experimental data, add background noise to better simulate experimental data
if (not expData):
    back = maximg*0.06
    images_mc = images_mc + back
    maximg = maximg + back
    minimg = minimg + back

# Construct sources array
img_list = []
pos_list = []

index = 0
if (expData):  
    sources=[]
    for i in sourcesRaw:
        idX = i//source_shape[0] 
        sourceY = ((idX)*stepSize[1]) + initPos[1]
        if (idX%2 == 0):
            sourceX = ((i%source_shape[0])*stepSize[0]) + initPos[0]
        else:
            sourceX = ((source_shape[0]-1 - (i%source_shape[0]) )*stepSize[0]) + initPos[0]
        sources.append([sourceX, sourceY])

    del(sourcesRaw)
    gc.collect()

    for i in images_mc:
        # Calculate source postion (usually the position of the maximum)

        # Suplied as input, just discretize
        maxp = [int(round(sources[index][0]*px_per_cm)),
                im_shape[1] - int(round(sources[index][1]*px_per_cm))]
        index += 1
        
        # No supplied, call a function
        #maxp2 = max_pos_fit(i)  

        if maxp is not None: 
            # For the images we could find the source position, normalize image,
            # and store position and normalized image
            img_norm = np.array(((i-minimg)/(maximg-minimg)), dtype="float16")
            img_list.append([maxp,img_norm])
            pos_list.append(maxp)
        del(maxp)

    del(images_mc)

else:
    for i in images_mc:

        # Suplied as input, just discretize
        maxp = [int(round(sources[index][0]*px_per_cm) + im_shape[0]/2),
                int(round(sources[index][1]*px_per_cm + im_shape[1]/2))]

        index += 1
        
        # No supplied, call a function
        #maxp2 = max_pos_fit(i)  
        
        if maxp is not None: 
            # For the images we could find the source position, normalize image,
            # and store position and normalized image
            img_norm = np.array(((i-minimg)/(maximg-minimg)), dtype="float16")
            img_list.append([maxp,img_norm])
            pos_list.append(maxp)
        del(maxp)

    del(images_mc)

del(pos_list)
gc.collect()

plt.imshow(np.asarray(img_list[0][1], dtype="float"))

#%%
# Sample run
importlib.reload(tomo)
rOffset = 2
mups = 8.9
zSlices = 15
cropSize = 2

met = "median"
maxR = 2.2
minR = 1.8
inBining = 4
rBinning = 4

baseFigPath = "../../TomosintesisArg/"

plt.style.use("default")

tomo_array = tomo.tomoMuaReconstructionV2(img_list, px_per_cm,
                                          minR, maxR, d, mups, tomo.shiftFuncDG, 
                                          zSlices, source_shape, cropSize, recMethod=met, percVal=5,
                                          mupsHomo=mups, muaHomo=0.07, mupsInc=6.7, 
                                          incSize=0.5, rBinning=rBinning)


# Z and X slices
z_slices = range(0,np.shape(tomo_array)[2], 4)
x_slices = [15, 28]#range(0,np.shape(tomo_array)[0], 5)
zslices_per_cm = d/(np.shape(tomo_array)[2]-1)
xslices_per_cm = (source_shape[1]*stepSize[0])/(np.shape(tomo_array)[0] - 1)

tomo_min = np.amin(tomo_array)
tomo_max = np.amax(tomo_array)


for x in x_slices:
    f = plt.figure(figsize = (16, 8), dpi = 120)
    #sf.set_title("X " + "{:.2f}".format(x * xslices_per_cm))
    plot = plt.imshow(tomo_array[x,:,:].T, vmin = tomo_min, vmax = tomo_max)
    #mpl.gca().invert_yaxis()
    plt.xlabel("y axis")
    plt.ylabel("z axis")
    plt.colorbar(plot)
    # f.savefig(baseFigPath + "muaTomo_x" + str(x) +  "_a" + str(ang) + "_r" 
            #    + str(roff) + "_" + met + ".png", transparent=True)


for z in z_slices:
    f = plt.figure(figsize = (16, 8), dpi=120)
    #f.set_title("Z " + "{:.2f}".format(z*zslices_per_cm))
    plot = plt.imshow(tomo_array[:, :, z], vmin = tomo_min, vmax=tomo_max)
    plt.gca().invert_yaxis()
    plt.xlabel("x axis")
    plt.ylabel("y axis")
    plt.colorbar(plot)
    # f.savefig(baseFigPath + "muaTomo_z" + str(z) +  "_a" + str(ang) + "_r" 
                # + str(roff) + "_" + met + ".png", transparent=True)



#%%
# Do the reconstruction for diferent sum methods
d=5
importlib.reload(tomo)
rOffset = 2
mups = 8.9
zSlices = 25
cropSize=2

# Set light plot style
plt.style.use("default")

angleStep = [45]
rOffsetMin = 1.8
rOffsetMax = 2.2
rOffsetNum = [4]
#rOffsetRange = [10]
methods = ["median"]
# methods = ["average","median","percentile"]


# muaHomo = 0.1
# muaInc = 0.2
# refFront1 = np.flip(256-imageio.imread("../Dibujos/Lab-VistaFrontal-1erInc.tif"),0)
# refFront1Norm = tomo.change_contrastMinMax(refFront1, muaHomo, muaInc)
# # f = plt.figure(figsize = (16, 8), dpi = 120)
# # plt.imshow(refFront1Norm)
# # plt.colorbar()

# refFront2 = np.flip(256-imageio.imread("../Dibujos/Lab-VistaFrontal-2doInc.tif"),0)
# refFront2Norm = tomo.change_contrastMinMax(refFront2, muaHomo, muaInc)

# refSup1 = 256-imageio.imread("../Dibujos/Lab-d5-VistaSup-1erInc.tif")
# refSup1Norm = tomo.change_contrastMinMax(refSup1, muaHomo, muaInc)
# f = plt.figure(figsize = (16, 8), dpi = 120)
# plt.imshow(refSup1Norm)
# plt.colorbar()


# refSup2 = 256-imageio.imread("../Dibujos/Lab-d5-VistaSup-2doInc.tif")
# refSup2Norm = tomo.change_contrastMinMax(refSup2, muaHomo, muaInc)



# SSIMs = []
# L1s = []

# PPTX Presentation imports
# Docs: https://python-pptx.readthedocs.io/en/latest/
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from pptx import enum

# Initialize the PPTX
# prs = Presentation()
# # First Slice
# slide = prs.slides.add_slide(prs.slide_layouts[0])
# title = slide.shapes.title
# title.text = "Tomosintesys algorithm test"

# Disable interactive plots
#plt.ioff()

for ang in angleStep:
    for roff in rOffsetNum:
        for met in methods:
            print(ang, roff, met)

            tomo_array, angles, offsets = tomo.tomoMuaReconstruction(img_list, px_per_cm,
                                          ang, rOffsetMin, rOffsetMax, d, mups, tomo.shiftFuncDG, 
                                          zSlices, source_shape, cropSize, recMethod=met, percVal=20,
                                          mupsHomo=mups, muaHomo=0.1, mupsInc=6.7, incSize=0.9, numRSteps=roff)


            # # First inclusion
            # imgFrontInc1 = tomo_array[:,:,18]
            # # f = plt.figure(figsize = (16, 8), dpi = 120)
            # # plt.imshow(imgFrontInc1)
            # SSIMF1 = compare_ssim(imgFrontInc1, refFront1Norm)
            # L1F1 = np.sum(abs(refFront1Norm-imgFrontInc1))/(np.size(imgFrontInc1))
            
            # imgSupInc1 = tomo_array[28,:,:].T
            # f = plt.figure(figsize = (16, 8), dpi = 120)
            # plt.title(str(ang)+" "+str(roff)+" "+met)
            # plt.imshow(imgSupInc1)
            
            # SSIMS1 = compare_ssim(imgSupInc1, refSup1Norm)
            # L1S1 = np.sum(abs(refSup1Norm-imgSupInc1))/(np.size(imgSupInc1))
            
            # # Second inclusion
            # imgFrontInc2 = tomo_array[:,:,6]
            # # f = plt.figure(figsize = (16, 8), dpi = 120)
            # # plt.imshow(imgInc2)
            # SSIMF2 = compare_ssim(imgFrontInc2, refFront2Norm)
            # L1F2 = np.sum(abs(refFront2Norm-imgFrontInc2))/(np.size(imgFrontInc2))
            
            # imgSupInc2 = tomo_array[17,:,:].T
            # # f = plt.figure(figsize = (16, 8), dpi = 120)
            # # plt.imshow(imgSupInc2)
            # SSIMS2 = compare_ssim(imgSupInc2, refSup2Norm)
            # L1S2 = np.sum(abs(refSup2Norm-imgSupInc2))/(np.size(imgSupInc2))
            
            # print(met, ang, SSIMF1, SSIMS1, SSIMF2, SSIMS2)
            # L1s.append((met, ang, L1F1, L1S1, L1F2, L1S2))
            # SSIMs.append((met, ang, SSIMF1, SSIMS1, SSIMF2, SSIMS2))
            
            # Z and X slices
            z_slices = range(0,np.shape(tomo_array)[2], 2)
            x_slices = [15, 28]#range(0,np.shape(tomo_array)[0], 5)
            zslices_per_cm = d/(np.shape(tomo_array)[2]-1)
            xslices_per_cm = (source_shape[1]*stepSize[0])/(np.shape(tomo_array)[0] - 1)

            tomo_min = np.amin(tomo_array)
            tomo_max = np.amax(tomo_array)
            
            for x in x_slices:
                f = plt.figure(figsize = (16, 8), dpi = 120)
                #sf.set_title("X " + "{:.2f}".format(x * xslices_per_cm))
                plot = plt.imshow(tomo_array[x,  :,  :].T, vmin = tomo_min, vmax = tomo_max)
                #mpl.gca().invert_yaxis()
                plt.xlabel("y axis")
                plt.ylabel("z axis")
                plt.colorbar(plot)
                # f.savefig(baseFigPath + "x" + str(x) +  "_a" + str(ang) + "_r" + str(roff) 
                #           + "_" + met + ".png", transparent=True)

           
            for z in z_slices:
                f = plt.figure(figsize = (16, 8), dpi=120)

                #f.set_title("Z " + "{:.2f}".format(z*zslices_per_cm))
                plot = plt.imshow(tomo_array[:, :, z], vmin = tomo_min, vmax=tomo_max)
                plt.gca().invert_yaxis()
                plt.xlabel("x axis")
                plt.ylabel("y axis")
                plt.colorbar(plot)
                # f.savefig(baseFigPath + "z" + str(z) +  "_a" + str(ang) + "_r" + str(roff) 
                #           + "_" + met + ".png", transparent=True)
            
            # slide = prs.slides.add_slide(prs.slide_layouts[1])
            # title = slide.shapes.title
            # title.text = "Angle step " + str(ang) + ", Radial Step " + str(roff) + " - " + met            
            # pic = slide.shapes.add_picture(baseFigPath + "z" + str(z_slices[0]) + "_a" 
            #                                + str(ang) + "_r" + str(roff) + "_" + met + ".png",
            #                                Inches(0.5), Inches(1.5), width=Inches(4))
            # pic = slide.shapes.add_picture(baseFigPath + "z" + str(z_slices[1]) + "_a" 
            #                                + str(ang) + "_r" + str(roff) + "_" + met + ".png",
            #                                Inches(2.5), Inches(1.5), width=Inches(4))
            # pic = slide.shapes.add_picture(baseFigPath + "z" + str(z_slices[2]) + "_a" 
            #                                + str(ang) + "_r" + str(roff) + "_" + met + ".png",
            #                                Inches(4.5), Inches(1.5), width=Inches(4))
            # pic = slide.shapes.add_picture(baseFigPath + "z" + str(z_slices[3]) + "_a" 
            #                                + str(ang) + "_r" + str(roff) + "_" + met + ".png",
            #                                Inches(0.5), Inches(3.5), width=Inches(4))
            # pic = slide.shapes.add_picture(baseFigPath + "z" + str(z_slices[4]) + "_a" 
            #                                + str(ang) + "_r" + str(roff) + "_" + met + ".png",
            #                                Inches(2.5), Inches(3.5), width=Inches(4))
            # pic = slide.shapes.add_picture(baseFigPath + "z" + str(z_slices[5]) + "_a" 
            #                                + str(ang) + "_r" + str(roff) + "_" + met + ".png",
            #                                Inches(4.5), Inches(3.5), width=Inches(4))
            # pic = slide.shapes.add_picture(baseFigPath + "x" + str(x_slices[0]) + "_a" 
            #                                + str(ang) + "_r" + str(roff) + "_" + met + ".png",
            #                                Inches(1), Inches(5.5), width=Inches(4))
            # pic = slide.shapes.add_picture(baseFigPath + "x" + str(x_slices[1]) + "_a" 
            #                                + str(ang) + "_r" + str(roff) + "_" + met + ".png",
            #                                Inches(5), Inches(5.5), width=Inches(4))

# prs.save(baseFigPath+"Presentation_Tomo2.pptx") 

#%%
# Different angle steps
angleStep = [22.5,45,90]
rOffsetRange = [10]
methods = ["median"]

# Disable interactive plots
plt.ioff()

for ang in angleStep:
    for roff in rOffsetRange:
        for met in methods:
            print(ang, roff, met)

            tomo_array, angles, offsets = tomo.tomoMuaReconstruction(img_list, px_per_cm,
                                          ang, rOffset, roff, d, mups, tomo.shiftFuncDG, 
                                          zSlices, source_shape, cropSize, recMethod=met, percVal=10,
                                          mupsHomo=mups, muaHomo=0.1, mupsInc=6.7, incSize=0.5)

            
            # Z and X slices
            z_slices = range(0,np.shape(tomo_array)[2], 2)
            x_slices = [15, 28]#range(0,np.shape(tomo_array)[0], 5)
            zslices_per_cm = d/(np.shape(tomo_array)[2]-1)
            xslices_per_cm = (source_shape[1]*stepSize[0])/(np.shape(tomo_array)[0] - 1)

            tomo_min = np.amin(tomo_array)
            tomo_max = np.amax(tomo_array)
            
            for x in x_slices:
                f = plt.figure(figsize = (16, 8), dpi = 120)
                #sf.set_title("X " + "{:.2f}".format(x * xslices_per_cm))
                plot = plt.imshow(tomo_array[x,  :,  :].T, vmin = tomo_min, vmax = tomo_max)
                #mpl.gca().invert_yaxis()
                plt.xlabel("y axis")
                plt.ylabel("z axis")
                plt.colorbar(plot)
                f.savefig(baseFigPath + "x" + str(x) +  "_a" + str(ang) + "_r" + str(roff) 
                          + "_" + met + ".png", transparent=True)
          
            for z in z_slices:
                f = plt.figure(figsize = (16, 8), dpi=120)

                #f.set_title("Z " + "{:.2f}".format(z*zslices_per_cm))
                plot = plt.imshow(tomo_array[:, :, z], vmin = tomo_min, vmax=tomo_max)
                plt.gca().invert_yaxis()
                plt.xlabel("x axis")
                plt.ylabel("y axis")
                plt.colorbar(plot)
                f.savefig(baseFigPath + "z" + str(z) +  "_a" + str(ang) + "_r" + str(roff) 
                          + "_" + met + ".png", transparent=True)
            
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            title.text = "Angle step " + str(ang) + ", Radial range " + str(roff) + " - " + met            
            pic = slide.shapes.add_picture(baseFigPath + "z" + str(z_slices[0]) + "_a" 
                                           + str(ang) + "_r" + str(roff) + "_" + met + ".png",
                                           Inches(0.5), Inches(1.5), width=Inches(4))
            pic = slide.shapes.add_picture(baseFigPath + "z" + str(z_slices[1]) + "_a" 
                                           + str(ang) + "_r" + str(roff) + "_" + met + ".png",
                                           Inches(2.5), Inches(1.5), width=Inches(4))
            pic = slide.shapes.add_picture(baseFigPath + "z" + str(z_slices[2]) + "_a" 
                                           + str(ang) + "_r" + str(roff) + "_" + met + ".png",
                                           Inches(4.5), Inches(1.5), width=Inches(4))
            pic = slide.shapes.add_picture(baseFigPath + "z" + str(z_slices[3]) + "_a" 
                                           + str(ang) + "_r" + str(roff) + "_" + met + ".png",
                                           Inches(0.5), Inches(3.5), width=Inches(4))
            pic = slide.shapes.add_picture(baseFigPath + "z" + str(z_slices[4]) + "_a" 
                                           + str(ang) + "_r" + str(roff) + "_" + met + ".png",
                                           Inches(2.5), Inches(3.5), width=Inches(4))
            pic = slide.shapes.add_picture(baseFigPath + "z" + str(z_slices[5]) + "_a" 
                                           + str(ang) + "_r" + str(roff) + "_" + met + ".png",
                                           Inches(4.5), Inches(3.5), width=Inches(4))
            pic = slide.shapes.add_picture(baseFigPath + "x" + str(x_slices[0]) + "_a" 
                                           + str(ang) + "_r" + str(roff) + "_" + met + ".png",
                                           Inches(1), Inches(5.5), width=Inches(4))
            pic = slide.shapes.add_picture(baseFigPath + "x" + str(x_slices[1]) + "_a" 
                                           + str(ang) + "_r" + str(roff) + "_" + met + ".png",
                                           Inches(5), Inches(5.5), width=Inches(4))

#%%
# Different range of radial offset

angleStep = [45]
rOffsetRange = [2,5,10]
methods = ["average"]

# Disable interactive plots
plt.ioff()

for ang in angleStep:
    for roff in rOffsetRange:
        for met in methods:
            print(ang, roff, met)

            tomo_array, angles, offsets = tomo.tomoMuaReconstruction(img_list, px_per_cm,
                                          ang, rOffset, roff, d, mups, tomo.shiftFuncDG, 
                                          zSlices, source_shape, cropSize, recMethod=met, percVal=10,
                                          mupsHomo=mups, muaHomo=0.07, mupsInc=6.7, incSize=0.5)

            
            # Z and X slices
            z_slices = range(0,np.shape(tomo_array)[2], 2)
            x_slices = [15, 28]#range(0,np.shape(tomo_array)[0], 5)
            zslices_per_cm = d/(np.shape(tomo_array)[2]-1)
            xslices_per_cm = (source_shape[1]*stepSize[0])/(np.shape(tomo_array)[0] - 1)

            tomo_min = np.amin(tomo_array)
            tomo_max = np.amax(tomo_array)

            
            for x in x_slices:
                f = plt.figure(figsize = (16, 8), dpi = 120)
                #sf.set_title("X " + "{:.2f}".format(x * xslices_per_cm))
                plot = plt.imshow(tomo_array[x,  :,  :].T, vmin = tomo_min, vmax = tomo_max)
                #mpl.gca().invert_yaxis()
                plt.xlabel("y axis")
                plt.ylabel("z axis")
                plt.colorbar(plot)
                f.savefig(baseFigPath + "x" + str(x) +  "_a" + str(ang) + "_r" + str(roff) 
                          + "_" + met + ".png", transparent=True)

           
            for z in z_slices:
                f = plt.figure(figsize = (16, 8), dpi=120)

                #f.set_title("Z " + "{:.2f}".format(z*zslices_per_cm))
                plot = plt.imshow(tomo_array[:, :, z], vmin = tomo_min, vmax=tomo_max)
                plt.gca().invert_yaxis()
                plt.xlabel("x axis")
                plt.ylabel("y axis")
                plt.colorbar(plot)
                f.savefig(baseFigPath + "z" + str(z) +  "_a" + str(ang) + "_r" + str(roff) + "_" + met + ".png", transparent=True)
            
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            title.text = "Angle step " + str(ang) + ", Radial range " + str(roff) + " - " + met            
            pic = slide.shapes.add_picture(baseFigPath + "z" + str(z_slices[0]) + "_a" 
                                           + str(ang) + "_r" + str(roff) + "_" + met + ".png",
                                           Inches(0.5), Inches(1.5), width=Inches(4))
            pic = slide.shapes.add_picture(baseFigPath + "z" + str(z_slices[1]) + "_a" 
                                           + str(ang) + "_r" + str(roff) + "_" + met + ".png",
                                           Inches(2.5), Inches(1.5), width=Inches(4))
            pic = slide.shapes.add_picture(baseFigPath + "z" + str(z_slices[2]) + "_a" 
                                           + str(ang) + "_r" + str(roff) + "_" + met + ".png",
                                           Inches(4.5), Inches(1.5), width=Inches(4))
            pic = slide.shapes.add_picture(baseFigPath + "z" + str(z_slices[3]) + "_a" 
                                           + str(ang) + "_r" + str(roff) + "_" + met + ".png",
                                           Inches(0.5), Inches(3.5), width=Inches(4))
            pic = slide.shapes.add_picture(baseFigPath + "z" + str(z_slices[4]) + "_a" 
                                           + str(ang) + "_r" + str(roff) + "_" + met + ".png",
                                           Inches(2.5), Inches(3.5), width=Inches(4))
            pic = slide.shapes.add_picture(baseFigPath + "z" + str(z_slices[5]) + "_a" 
                                           + str(ang) + "_r" + str(roff) + "_" + met + ".png",
                                           Inches(4.5), Inches(3.5), width=Inches(4))
            pic = slide.shapes.add_picture(baseFigPath + "x" + str(x_slices[0]) + "_a" 
                                           + str(ang) + "_r" + str(roff) + "_" + met + ".png",
                                           Inches(1), Inches(5.5), width=Inches(4))
            pic = slide.shapes.add_picture(baseFigPath + "x" + str(x_slices[1]) + "_a" 
                                           + str(ang) + "_r" + str(roff) + "_" + met + ".png",
                                           Inches(5), Inches(5.5), width=Inches(4))
#%%
prs.save(baseFigPath+"Presentation_Tomo.pptx")  

#%%
importlib.reload(tomo)

# %%
print(np.max(refFront1), np.min(refFront1))
print(np.max(refFront1Norm), np.min(refFront1Norm))

# %%
